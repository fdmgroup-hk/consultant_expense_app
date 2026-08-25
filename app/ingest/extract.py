"""Pull readable text out of consultant material.

Each extractor yields ``Segment`` records that keep a human-readable locator
("Slide 7", "Page 3") so answers can cite exactly where a fact came from.
"""
from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

SUPPORTED_EXTENSIONS = {".pptx", ".pdf", ".docx", ".md", ".markdown", ".txt"}


@dataclass
class Segment:
    locator: str
    heading: str
    text: str


class UnsupportedFileType(ValueError):
    pass


def _clean(text: str) -> str:
    text = text.replace("\x00", " ").replace("\u00a0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _extract_pptx(path: Path) -> list[Segment]:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        deck = Presentation(str(path))
    except PackageNotFoundError as exc:  # .ppt saved with a .pptx name, or corrupt
        raise UnsupportedFileType(
            "Could not open as PowerPoint. If this is a legacy .ppt file, "
            "re-save it as .pptx and upload again."
        ) from exc

    segments: list[Segment] = []
    for index, slide in enumerate(deck.slides, start=1):
        title = ""
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
        except (AttributeError, ValueError):
            title = ""

        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                body = shape.text_frame.text.strip()
                if body and body != title:
                    parts.append(body)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

        # Speaker notes are where consultants explain what the slide actually means.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")

        body = _clean("\n".join(parts))
        if not body and not title:
            continue
        segments.append(
            Segment(locator=f"Slide {index}", heading=title or f"Slide {index}", text=body)
        )
    return segments


def _extract_pdf(path: Path) -> list[Segment]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    segments: list[Segment] = []
    for index, page in enumerate(reader.pages, start=1):
        body = _clean(page.extract_text() or "")
        if not body:
            continue
        first_line = body.splitlines()[0][:120] if body.splitlines() else ""
        segments.append(
            Segment(locator=f"Page {index}", heading=first_line, text=body)
        )
    return segments


def _extract_docx(path: Path) -> list[Segment]:
    from docx import Document

    doc = Document(str(path))
    segments: list[Segment] = []
    heading = ""
    buffer: list[str] = []
    counter = 0

    def flush() -> None:
        nonlocal buffer, counter
        body = _clean("\n".join(buffer))
        if body:
            counter += 1
            segments.append(
                Segment(
                    locator=f"Section {counter}",
                    heading=heading or f"Section {counter}",
                    text=body,
                )
            )
        buffer = []

    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("heading") or style == "title":
            flush()
            heading = text
        else:
            buffer.append(text)
    flush()

    for t_index, table in enumerate(doc.tables, start=1):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in table.rows
            if any(cell.text.strip() for cell in row.cells)
        ]
        if rows:
            segments.append(
                Segment(
                    locator=f"Table {t_index}",
                    heading=f"Table {t_index}",
                    text=_clean("\n".join(rows)),
                )
            )
    return segments


def _extract_markdown(path: Path) -> list[Segment]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    segments: list[Segment] = []
    heading = ""
    buffer: list[str] = []
    counter = 0

    def flush() -> None:
        nonlocal buffer, counter
        body = _clean("\n".join(buffer))
        if body:
            counter += 1
            segments.append(
                Segment(
                    locator=heading or f"Section {counter}",
                    heading=heading or f"Section {counter}",
                    text=body,
                )
            )
        buffer = []

    for line in raw.splitlines():
        match = re.match(r"^(#{1,6})\s+(.*)$", line)
        if match:
            flush()
            heading = match.group(2).strip()
        else:
            buffer.append(line)
    flush()
    return segments


def _extract_txt(path: Path) -> list[Segment]:
    body = _clean(path.read_text(encoding="utf-8", errors="replace"))
    return [Segment(locator="Document", heading="", text=body)] if body else []


_EXTRACTORS = {
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".md": _extract_markdown,
    ".markdown": _extract_markdown,
    ".txt": _extract_txt,
}


def extract(path: Path) -> list[Segment]:
    suffix = path.suffix.lower()
    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise UnsupportedFileType(
            f"{suffix or 'This file type'} is not supported. "
            f"Upload one of: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    return extractor(path)
